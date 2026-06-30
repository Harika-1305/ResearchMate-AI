from pptx import Presentation

def create_ppt(content):

    prs = Presentation()

    slides_data = content.split("Slide ")

    for slide_data in slides_data:

        slide_data = slide_data.strip()

        if not slide_data:
            continue

        slide = prs.slides.add_slide(
            prs.slide_layouts[1]
        )

        lines = slide_data.split("\n")

        title = lines[0]

        slide.shapes.title.text = title

        body = slide.placeholders[1]

        body.text = "\n".join(lines[1:])

    ppt_path = "ResearchMate_Presentation.pptx"

    prs.save(ppt_path)

    return ppt_path